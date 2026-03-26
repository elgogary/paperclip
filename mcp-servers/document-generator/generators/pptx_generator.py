"""PPTX generator — populates a branded .pptx template with slide content."""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from .diagram_renderer import render_mermaid


def _hex_to_rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _add_title_bar(slide, prs, title_text: str, brand: dict):
    """Add a colored title bar at the top of the slide."""
    primary = _hex_to_rgb(brand.get("primary", "#007bff"))
    w, h = prs.slide_width, Inches(1.1)
    shape = slide.shapes.add_shape(1, 0, 0, w, h)  # MSO_SHAPE.RECTANGLE
    shape.fill.solid()
    shape.fill.fore_color.rgb = primary
    shape.line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.5)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.font.bold = True
    p.font.name = brand.get("font", "Calibri")


def _add_body_text(slide, prs, top: float, items: list, brand: dict, font_size: int = 18):
    """Add bullet points or text below the title bar."""
    left, width = Inches(0.6), prs.slide_width - Inches(1.2)
    height = prs.slide_height - Inches(top + 0.3)
    txBox = slide.shapes.add_textbox(left, Inches(top), width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = _hex_to_rgb(brand.get("text", "#333333"))
        p.font.name = brand.get("font", "Calibri")
        p.space_after = Pt(8)


def _slide_title(slide, prs, data: dict, brand: dict):
    """Full-color title slide."""
    primary = _hex_to_rgb(brand.get("primary", "#007bff"))
    shape = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = primary
    shape.line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(1)
    tf.margin_top = Inches(2)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    p = tf.paragraphs[0]
    p.text = data.get("title", "")
    p.font.size = Pt(36)
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.font.bold = True
    p.font.name = brand.get("font", "Calibri")
    p.alignment = PP_ALIGN.LEFT

    if data.get("subtitle"):
        p2 = tf.add_paragraph()
        p2.text = data["subtitle"]
        p2.font.size = Pt(18)
        p2.font.color.rgb = RGBColor(0xDD, 0xDD, 0xFF)
        p2.font.name = brand.get("font", "Calibri")
        p2.space_before = Pt(20)


def _slide_section(slide, prs, data: dict, brand: dict):
    """Section divider — full color background."""
    _slide_title(slide, prs, data, brand)


def _slide_content(slide, prs, data: dict, brand: dict):
    """Content slide with title bar + bullets."""
    _add_title_bar(slide, prs, data.get("title", ""), brand)
    bullets = data.get("bullets", [])
    if bullets:
        _add_body_text(slide, prs, 1.3, bullets, brand)


def _slide_table(slide, prs, data: dict, brand: dict):
    """Table slide with title bar + table."""
    _add_title_bar(slide, prs, data.get("title", ""), brand)

    tbl_data = data.get("table", {})
    headers = tbl_data.get("headers", [])
    rows_data = tbl_data.get("rows", [])
    if not headers:
        return

    cols = len(headers)
    total_rows = 1 + len(rows_data)
    left, top = Inches(0.5), Inches(1.4)
    width = prs.slide_width - Inches(1.0)
    height = Inches(min(total_rows * 0.4, 5.0))

    table = slide.shapes.add_table(total_rows, cols, left, top, width, height).table

    # Header row
    primary = _hex_to_rgb(brand.get("primary", "#007bff"))
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            p.font.name = brand.get("font", "Calibri")
        cell.fill.solid()
        cell.fill.fore_color.rgb = primary

    # Data rows
    for r, row_data in enumerate(rows_data):
        for c, val in enumerate(row_data):
            if c < cols:
                cell = table.cell(r + 1, c)
                cell.text = str(val)
                for p in cell.text_frame.paragraphs:
                    p.font.size = Pt(12)
                    p.font.name = brand.get("font", "Calibri")
                    p.font.color.rgb = _hex_to_rgb(brand.get("text", "#333333"))
                # Alternate row colors
                if r % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(0xF0, 0xF4, 0xF8)


def _slide_diagram(slide, prs, data: dict, brand: dict, output_dir: Path):
    """Diagram slide — renders mermaid to PNG and embeds it."""
    _add_title_bar(slide, prs, data.get("title", ""), brand)

    mermaid_code = data.get("mermaid", "")
    if mermaid_code:
        slug = data.get("title", "diagram").lower().replace(" ", "_")[:30]
        img_path = output_dir / f"_diagram_{slug}.png"
        try:
            render_mermaid(mermaid_code, img_path, width=1000, height=600)
            if img_path.exists():
                slide.shapes.add_picture(
                    str(img_path), Inches(0.8), Inches(1.5),
                    width=prs.slide_width - Inches(1.6),
                )
                if data.get("caption"):
                    _add_body_text(slide, prs, 6.2, [data["caption"]], brand, font_size=12)
                return
        except Exception:
            pass

    # Fallback: show mermaid code as text
    _add_body_text(slide, prs, 1.3, [f"[Diagram: {mermaid_code[:200]}]"], brand, font_size=14)
    if data.get("caption"):
        _add_body_text(slide, prs, 5.5, [data["caption"]], brand, font_size=12)


def _slide_image(slide, prs, data: dict, brand: dict):
    """Image slide with title and optional caption."""
    _add_title_bar(slide, prs, data.get("title", ""), brand)

    img_path = data.get("image", "")
    if img_path and Path(img_path).exists():
        slide.shapes.add_picture(
            img_path, Inches(0.8), Inches(1.5),
            width=prs.slide_width - Inches(1.6),
        )
    else:
        _add_body_text(slide, prs, 2.5, [f"[Image placeholder: {img_path}]"], brand)

    if data.get("caption"):
        _add_body_text(slide, prs, 6.2, [data["caption"]], brand, font_size=12)


LAYOUT_HANDLERS = {
    "title": _slide_title,
    "section": _slide_section,
    "content": _slide_content,
    "table": _slide_table,
    "image": _slide_image,
}


def generate_pptx(template_path: Path | None, data: dict, output_path: Path) -> dict:
    """Generate a .pptx from template + slide data.

    Returns dict with metadata about the generated presentation.
    """
    brand = data.get("brand", {"primary": "#007bff", "accent": "#28a745", "text": "#333333"})

    if template_path and template_path.exists():
        prs = Presentation(str(template_path))
    else:
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

    slides_data = data.get("slides", [])
    diagram_count = 0

    for slide_data in slides_data:
        layout = slide_data.get("layout", "content")
        blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
        slide = prs.slides.add_slide(blank_layout)

        if layout == "diagram":
            _slide_diagram(slide, prs, slide_data, brand, output_path.parent)
            diagram_count += 1
        elif layout in LAYOUT_HANDLERS:
            LAYOUT_HANDLERS[layout](slide, prs, slide_data, brand)
        else:
            _slide_content(slide, prs, slide_data, brand)

        # Speaker notes
        if slide_data.get("notes"):
            slide.notes_slide.notes_text_frame.text = slide_data["notes"]

    prs.save(str(output_path))

    return {
        "slides": len(slides_data),
        "diagrams_rendered": diagram_count,
    }
